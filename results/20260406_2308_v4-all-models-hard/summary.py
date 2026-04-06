#!/usr/bin/env python3
"""Generate a 2-slide PPTX summary for the v4 all-models-hard batch."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pathlib import Path

# ---------------------------------------------------------------------------
# Data
# ---------------------------------------------------------------------------

TITLE = "v4 Cheatsheet Evaluation — All Models x Hard Datasets"
SUBTITLE = "SAIR Mathematics Distillation Challenge  |  2026-04-06  |  Prompt: v4_10KB_cheatsheet.md"

MODELS = [
    # model, (correct, scorable, parse_err) per dataset, bias
    # hard1=69, hard2=200, hard3=400
    ("grok-4.1-fast",     (53,68,1),   (153,200,0), (242,400,0), "Balanced"),
    ("gemma-4-31b",       (43,69,0),   (145,197,3), (202,385,15),"TRUE bias (confab)"),
    ("gemini-flash-lite", (51,69,0),   (104,200,0), (225,400,0), "FALSE bias"),
    ("deepseek-v3.2",    (45,69,0),   (101,200,0), (206,400,0), "Total FALSE (0% TRUE recall)"),
    ("llama-3.3-70b",    (45,69,0),   (100,200,0), (206,400,0), "Total FALSE (0% TRUE recall)"),
    ("gpt-oss-20b",      (38,65,4),   (76,174,26), (179,358,42),"TRUE bias + parse err"),
    ("gpt-oss-120b",     (29,66,3),   (102,194,6), (170,369,31),"Extreme TRUE bias"),
]

HEADERS = ["Model", "hard1 (69)", "hard2 (200)", "hard3 (400)", "Parse Err", "Bias Pattern"]

# Colours
BG        = RGBColor(0x1A, 0x1A, 0x2E)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT     = RGBColor(0xCC, 0xCC, 0xCC)
ACCENT    = RGBColor(0x4E, 0xC9, 0xB0)
GOLD      = RGBColor(0xFF, 0xD7, 0x00)
RED       = RGBColor(0xFF, 0x6B, 0x6B)
ORANGE    = RGBColor(0xFF, 0xA5, 0x00)
HEADER_BG = RGBColor(0x2D, 0x2D, 0x4E)
ROW_EVEN  = RGBColor(0x22, 0x22, 0x3A)
ROW_ODD   = RGBColor(0x1A, 0x1A, 0x2E)
GREEN_BG  = RGBColor(0x1B, 0x3A, 0x2A)
RED_BG    = RGBColor(0x3A, 0x1B, 0x1B)

# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text_box(slide, left, top, width, height, text, font_size=12,
                 color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def set_cell(cell, text, font_size=10, color=WHITE, bold=False, align=PP_ALIGN.CENTER, bg_color=None):
    cell.text = ""
    p = cell.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = "Calibri"
    p.alignment = align
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    # margins
    cell.margin_left = Emu(45720)
    cell.margin_right = Emu(45720)
    cell.margin_top = Emu(27432)
    cell.margin_bottom = Emu(27432)
    if bg_color:
        cell_fill = cell.fill
        cell_fill.solid()
        cell_fill.fore_color.rgb = bg_color


def acc_color(acc_str):
    val = float(acc_str.rstrip('%'))
    if val >= 70: return ACCENT
    if val >= 55: return GOLD
    return RED


def bias_color(bias):
    if "Balanced" in bias: return ACCENT
    if "FALSE" in bias: return LIGHT
    return RED


# ---------------------------------------------------------------------------
# Slide 1 — Results Table
# ---------------------------------------------------------------------------

def build_slide1(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(slide, BG)

    # Title
    add_text_box(slide, Inches(0.4), Inches(0.2), Inches(9.2), Inches(0.5),
                 TITLE, font_size=22, color=WHITE, bold=True)
    add_text_box(slide, Inches(0.4), Inches(0.65), Inches(9.2), Inches(0.3),
                 SUBTITLE, font_size=10, color=LIGHT)

    # Table
    rows = len(MODELS) + 1
    cols = 6
    tbl_left = Inches(0.2)
    tbl_top = Inches(1.15)
    tbl_width = Inches(9.6)
    tbl_height = Inches(3.2)
    shape = slide.shapes.add_table(rows, cols, tbl_left, tbl_top, tbl_width, tbl_height)
    table = shape.table

    # Column widths
    col_widths = [Inches(1.8), Inches(1.6), Inches(1.6), Inches(1.6), Inches(0.9), Inches(2.1)]
    for i, w in enumerate(col_widths):
        table.columns[i].width = w

    # Header row
    for j, hdr in enumerate(HEADERS):
        set_cell(table.cell(0, j), hdr, font_size=10, color=WHITE, bold=True, bg_color=HEADER_BG)

    # Helper to format a (correct, scorable, parse_err) tuple
    def fmt_acc(tup):
        correct, scorable, _ = tup
        pct = correct / scorable * 100 if scorable > 0 else 0
        return f"{correct}/{scorable} ({pct:.1f}%)", pct

    # Data rows
    for i, (model, h1, h2, h3, bias) in enumerate(MODELS):
        row_bg = ROW_EVEN if i % 2 == 0 else ROW_ODD
        row_idx = i + 1
        set_cell(table.cell(row_idx, 0), model, font_size=9, color=WHITE, bold=True, align=PP_ALIGN.LEFT, bg_color=row_bg)
        for col_idx, ds in enumerate((h1, h2, h3), start=1):
            label, pct = fmt_acc(ds)
            set_cell(table.cell(row_idx, col_idx), label, font_size=9, color=acc_color(f"{pct:.1f}%"), bold=True, bg_color=row_bg)
        total_parse = h1[2] + h2[2] + h3[2]
        pe_color = RED if total_parse > 20 else (GOLD if total_parse > 0 else ACCENT)
        set_cell(table.cell(row_idx, 4), str(total_parse), font_size=10, color=pe_color, bold=True, bg_color=row_bg)
        set_cell(table.cell(row_idx, 5), bias, font_size=8, color=bias_color(bias), bg_color=row_bg)

    # Highlight grok row
    for j in range(cols):
        cell = table.cell(1, j)
        cell_fill = cell.fill
        cell_fill.solid()
        cell_fill.fore_color.rgb = GREEN_BG

    # Footer stats
    add_text_box(slide, Inches(0.4), Inches(4.55), Inches(9.2), Inches(0.7),
                 "4,683 API calls  |  669 problems/model  |  Total cost: ~$2.50  |  Datasets: hard1 (69), hard2 (200), hard3 (400)",
                 font_size=9, color=LIGHT, alignment=PP_ALIGN.CENTER)

    # Legend
    legend_y = Inches(4.95)
    add_text_box(slide, Inches(0.4), legend_y, Inches(9.2), Inches(0.3),
                 "Color key:  >=70% (green)  |  55-69% (gold)  |  <55% (red)        "
                 "Bias:  Balanced (green)  |  FALSE bias (grey)  |  TRUE bias (red)",
                 font_size=8, color=LIGHT, alignment=PP_ALIGN.CENTER)


# ---------------------------------------------------------------------------
# Slide 2 — Insights & Recommendations
# ---------------------------------------------------------------------------

INSIGHTS = [
    ("1024-token cap is the #1 GPT-OSS failure driver",
     "GPT-OSS-120b: ~90% of failures hit exactly 1024 tokens mid-reasoning. "
     "No VERDICT emitted; harness defaults unpredictably. Community best practice: emit VERDICT as first output line, reason after."),
    ("v4 creates extreme polarization across model families",
     "DeepSeek/Llama/Gemini: 0-9% TRUE recall (always FALSE). GPT-OSS: 85-91% TRUE recall (always TRUE). "
     "Only Grok stays balanced. The step-by-step decision procedure is too rigid for weaker models and too loose for stronger ones."),
    ("Betka's 98% sheet (community top) uses contradiction motifs C1-C14",
     "Our v4 motifs over-trigger on FALSE problems (59 confabs in Grok alone). "
     "Betka's sheet forces feature computation (rhsVars, topShape, Lx/Rx, xTop) before any rule fires, with hard-stop on first match."),
    ("Llama/DeepSeek treat cheatsheet as pure lookup table",
     "317/318 Llama failures are identical 'RULE: default false' boilerplate (median 88 tokens). "
     "No reasoning attempted. DeepSeek fabricates counterexamples (38%) when no rule matches."),
    ("Gemini hallucinates structural rules not in the cheatsheet",
     "Invents 'rightmost exclusion' and 'count exclusion' heuristics (53% of failures). "
     "Pattern-matches the cheatsheet's naming convention to fabricate plausible but invalid rules."),
]

RECOMMENDATIONS = [
    ("VERDICT-first output format", "critical",
     "Emit 'VERDICT: TRUE' or 'VERDICT: FALSE' as the first line, then reasoning. "
     "Eliminates all token-truncation parse errors (72 on GPT-20b, 40 on GPT-120b). "
     "Community consensus (Betka, Heath, stokarz)."),
    ("Counterexample hard-stop with named rule lock", "critical",
     "Once a counterexample is found: cite rule, output FALSE immediately, stop reasoning. "
     "GPT-120b finds correct counterexamples then 'talks itself out of it' with fluent wrong proofs. "
     "Betka/stokarz both use this. Prevents 288 GPT-120b confabulations."),
    ("Add Spine Isolation Theorem (McKenna, EQT01-000011)", "high",
     "One sentence: 'A pure left-spine equation can only imply equations that are themselves pure left-spine of equal or greater depth.' "
     "Validated on 1.54M pairs with zero exceptions. Fixes +32 hard3 problems across all models."),
    ("Adopt Betka's feature-first protocol for TRUE routes", "high",
     "Compute source features (rhsVars, topShape, Lx, Rx, xTop, square) BEFORE checking motifs. "
     "Our v4 lets models skip feature computation and pattern-match rule names instead. "
     "Betka achieves 98% on hard200 with GPT-120b using this approach."),
    ("Force minimal reasoning before default-FALSE", "medium",
     "Add: 'Before defaulting to FALSE, attempt one substitution x=y=z and check if Eq2 follows from Eq1.' "
     "Llama outputs 'default false' in 88 tokens with zero analysis. A forced check would catch trivial TRUE cases."),
    ("Increase max_tokens to 2048 for GPT-OSS models", "medium",
     "GPT-120b correct answers average 1063-1188 tokens. Current 1024 cap truncates ~90% of failures. "
     "Doubling to 2048 costs ~$0.10 more per run but recovers most truncated verdicts."),
]


PRIO_COLORS = {
    "critical": RGBColor(0xFF, 0x4D, 0x4D),
    "high":     RGBColor(0xFF, 0xA5, 0x00),
    "medium":   RGBColor(0xFF, 0xD7, 0x00),
}


def build_slide2(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG)

    # Title
    add_text_box(slide, Inches(0.4), Inches(0.15), Inches(9.2), Inches(0.4),
                 "Insights & Recommendations", font_size=20, color=WHITE, bold=True)

    # --- Left column: Insights ---
    col_left = Inches(0.3)
    col_width = Inches(4.5)
    y = Inches(0.55)

    add_text_box(slide, col_left, y, col_width, Inches(0.25),
                 "KEY INSIGHTS (from JSON analysis + community intel)", font_size=9, color=ACCENT, bold=True)
    y += Inches(0.28)

    for title, detail in INSIGHTS:
        add_text_box(slide, col_left, y, col_width, Inches(0.2),
                     ">  " + title, font_size=8, color=GOLD, bold=True)
        y += Inches(0.2)
        add_text_box(slide, Inches(0.5), y, Inches(4.3), Inches(0.45),
                     detail, font_size=7, color=LIGHT)
        y += Inches(0.45)

    # --- Right column: Recommendations ---
    col_right = Inches(5.0)
    col_width_r = Inches(4.8)
    y2 = Inches(0.55)

    add_text_box(slide, col_right, y2, col_width_r, Inches(0.25),
                 "NEXT STEPS (prioritized)", font_size=9, color=ACCENT, bold=True)
    y2 += Inches(0.32)

    for i, (title, priority, detail) in enumerate(RECOMMENDATIONS, 1):
        prio_color = PRIO_COLORS.get(priority, LIGHT)
        # Priority tag + title
        tag = f"[{priority.upper()}]"
        add_text_box(slide, col_right, y2, Inches(0.8), Inches(0.18),
                     tag, font_size=7, color=prio_color, bold=True)
        add_text_box(slide, Inches(5.8), y2, Inches(4.0), Inches(0.18),
                     f"{i}. {title}", font_size=8, color=WHITE, bold=True)
        y2 += Inches(0.19)
        add_text_box(slide, Inches(5.8), y2, Inches(4.0), Inches(0.5),
                     detail, font_size=7, color=LIGHT)
        y2 += Inches(0.5)

    # Sources footer
    add_text_box(slide, Inches(0.3), Inches(5.05), Inches(9.4), Inches(0.2),
                 "Sources: JSON reasoning analysis  |  Betka (EQT01-000014, 98% hard200)  |  McKenna (Spine Isolation, +32 hard3)  |  "
                 "Tao (50/50 eval confirmed)  |  stokarz (199/200 hard200)",
                 font_size=6, color=LIGHT, alignment=PP_ALIGN.CENTER)

    # Footer
    add_text_box(slide, Inches(0.3), Inches(5.25), Inches(9.4), Inches(0.2),
                 "SAIR Mathematics Distillation Challenge -- Equational Theories Stage 1  |  Deadline: 2026-04-20",
                 font_size=7, color=LIGHT, alignment=PP_ALIGN.CENTER)


# ---------------------------------------------------------------------------
# Appendix — one slide per model with full confusion matrix
# ---------------------------------------------------------------------------

# Per-model, per-dataset: (accuracy_str, TT_str, TF, FT, FF, parse_err)
# TT_str = "correct/total_true" to show recall denominator
MODEL_DETAILS = [
    {
        "name": "grok-4.1-fast",
        "display": "Grok 4.1 Fast",
        "datasets": [
            ("hard1",  "77.9%", "16/24",  8,   7,  37, 1),
            ("hard2",  "76.5%", "75/100", 25,  22, 78, 0),
            ("hard3",  "60.5%", "102/195",93,  65, 140,0),
        ],
        "insight": "Most balanced model -- 67% TRUE recall on hard1, 75% on hard2. Only model above 60% on hard3.",
        "failures": [
            ("Default FALSE fallback (no TRUE route fired)", "126 (57%)",
             "Cheatsheet TRUE-detection rules too narrow; model correctly follows protocol but defaults to FALSE on unrecognized TRUE patterns."),
            ("Contradiction motif over-trigger", "59 (27%)",
             "C8/C9/C13 motifs (rhsVars, Lx/Rx flags, topShape) fire on FALSE problems. Motifs derived from biased sample; structural match doesn't guarantee implication."),
            ("Constant-operation rule misfire", "19 (9%)",
             "Disjoint variable sets on Eq1 sides wrongly triggers 'constant magma' conclusion. Rule is mathematically overbroad."),
            ("Projection rule misidentification", "10 (5%)",
             "Left/right projection rule fires but T actually contains x, or derivation from projection to Eq2 is incorrect."),
        ],
    },
    {
        "name": "gemma-4-31b",
        "display": "Gemma 4 31B",
        "datasets": [
            ("hard1",  "62.3%", "21/24",   3,  23, 22,  0),
            ("hard2",  "73.6%", "85/97",  12,  40, 60,  3),
            ("hard3",  "52.5%", "149/189",40, 143, 53, 15),
        ],
        "insight": "Strong TRUE recall (88% on hard1/hard2) but heavy confabulation on FALSE problems. 18 parse errors (empty API responses).",
        "failures": [
            ("Contradiction motif misfire (FALSE\u2192TRUE)", "~153 (55%)",
             "Model walks through motif checklist, fails to match, then reverses mid-reasoning to output TRUE. Confabulation at the rule-application layer."),
            ("Left projection misfire (FALSE\u2192TRUE)", "~43 (15%)",
             "Claims 'left projection' rule applies but structural conditions don't hold. Pattern-matches rule name without verifying preconditions."),
            ("Right-zero witness misuse (TRUE\u2192FALSE)", "35 (13%)",
             "When Eq1 fails in right-zero magma, concludes FALSE. Logic error: a magma where Eq1 fails is NOT a valid counterexample."),
            ("Syntactic exclusion misfires (TRUE\u2192FALSE)", "15 (5%)",
             "Leftmost/count exclusion rules applied without verifying preconditions. Pattern-matching syntax without understanding semantics."),
            ("Parse errors (empty API response)", "18 (6%)",
             "output_tokens=0, empty content. API-level failures/timeouts, not model reasoning issues."),
        ],
    },
    {
        "name": "gemini-flash-lite",
        "display": "Gemini 3.1 Flash Lite",
        "datasets": [
            ("hard1",  "73.9%", "9/24",    15,   3, 42, 0),
            ("hard2",  "52.0%", "9/100",   91,   5, 95, 0),
            ("hard3",  "56.2%", "36/195", 159,  16,189, 0),
        ],
        "insight": "Massive FALSE bias -- 9% TRUE recall on hard2. Near-perfect FALSE accuracy but misses almost all TRUE implications. Zero parse errors.",
        "failures": [
            ("Fabricated rightmost exclusion heuristic", "154 (53%)",
             "Model invents a 'rightmost variable must match' rule not in the cheatsheet. Fires on ~58% of TRUE problems. Hallucinated structural heuristic."),
            ("Default FALSE on non-bare equations", "64 (22%)",
             "When Eq1 is not 'x = T' form, immediately concludes no TRUE routes apply. Skips all non-trivial reasoning."),
            ("Fabricated leftmost exclusion heuristic", "35 (12%)",
             "Symmetric to rightmost: checks leftmost variable preservation and rejects TRUE implications. Another hallucinated rule."),
            ("Contradiction motif over-firing (FALSE\u2192TRUE)", "15 (5%)",
             "Incorrectly claims equation 'collapses' magma. Accounts for 63% of the 24 false positives."),
            ("Count exclusion misapplication", "12 (4%)",
             "Claims variable occurrence counts must be preserved across implication. Not a valid algebraic invariant."),
        ],
    },
    {
        "name": "deepseek-v3.2",
        "display": "DeepSeek V3.2",
        "datasets": [
            ("hard1",  "65.2%", "0/24",   24,  0,  45, 0),
            ("hard2",  "50.5%", "1/100",  99,  0, 100, 0),
            ("hard3",  "51.5%", "2/195", 193,  1, 204, 0),
        ],
        "insight": "Total FALSE bias -- 0-1% TRUE recall. Perfect FALSE detection but no signal on TRUE. Zero parse errors.",
        "failures": [
            ("Step 6 'default FALSE' fallback", "165 (52%)",
             "No TRUE route fires; model follows cheatsheet correctly and defaults to FALSE. TRUE-detection rules have zero coverage on hard problems."),
            ("Fabricated / incorrect counterexamples", "121 (38%)",
             "Constructs small magma, claims Eq1 holds and Eq2 fails, but math is wrong -- Eq1 often fails in that magma too. Not a valid counterexample."),
            ("Wrong projection rule application", "~20 (6%)",
             "Correctly identifies left-zero/right-zero behavior for Eq1, then mis-reduces Eq2 with flawed variable trace algebra."),
            ("Wrong parity/syntactic exclusion", "~10 (3%)",
             "Applies parity exclusion rule with incorrect parity computation. 'Left side has one x, right has two' -- miscounted."),
        ],
    },
    {
        "name": "llama-3.3-70b",
        "display": "Llama 3.3 70B",
        "datasets": [
            ("hard1",  "65.2%", "0/24",   24,  0,  45, 0),
            ("hard2",  "50.0%", "0/100", 100,  0, 100, 0),
            ("hard3",  "51.5%", "2/195", 193,  1, 204, 0),
        ],
        "insight": "Total FALSE bias -- 0% TRUE recall on hard1/hard2. Treats cheatsheet as pure lookup table with no generalization.",
        "failures": [
            ("'RULE: default false' on every unmatched problem", "317 (99.7%)",
             "Identical boilerplate on all failures: 'The given equations do not match any exact TRUE routes.' Median 88 tokens output. Zero algebraic reasoning attempted."),
            ("No chain-of-thought reasoning", "317 (100%)",
             "Every response is a direct verdict + one-sentence justification. Never attempts substitution, variable analysis, or counterexample construction."),
            ("Only 2 TRUE problems matched (left projection)", "2 correct",
             "hard3_0002 and hard3_0005 matched the 'left projection' named rule exactly. Everything else falls through to default FALSE."),
        ],
    },
    {
        "name": "gpt-oss-20b",
        "display": "GPT-OSS 20B",
        "datasets": [
            ("hard1",  "58.5%", "16/23",   7, 20, 22,  4),
            ("hard2",  "43.7%", "53/89",  36, 62, 23, 26),
            ("hard3",  "50.0%", "130/174",44,135, 49, 42),
        ],
        "insight": "TRUE bias + 72 parse errors. Singleton collapse misfire is the dominant confabulation source.",
        "failures": [
            ("Singleton/collapse rule misfire (FALSE\u2192TRUE)", "217 (58%)",
             "Incorrectly triggers 'singleton collapse' heuristic. 183/217 contain the word 'singleton'. Fires on bare-variable left sides that superficially resemble TRUE patterns."),
            ("Token limit truncation (parse errors)", "72 (19%)",
             "All parse errors hit exactly 1024 output tokens. Model cut off mid-reasoning before emitting VERDICT line. Hard limit issue."),
            ("Parity exclusion misfire (TRUE\u2192FALSE)", "74 (20%)",
             "Named 'parity exclusion' fires incorrectly. Model computes parity mismatch and concludes FALSE on genuinely TRUE implications."),
            ("Fabricated counterexample (TRUE\u2192FALSE)", "31 (8%)",
             "Constructs left-zero or 2-element magma, checks Eq2 fails, declares FALSE -- but either the check is wrong or Eq1 doesn't hold in that magma."),
        ],
    },
    {
        "name": "gpt-oss-120b",
        "display": "GPT-OSS 120B",
        "datasets": [
            ("hard1",  "43.9%", "21/23",   2, 35,  8,  3),
            ("hard2",  "52.6%", "87/99",  12, 80, 15,  6),
            ("hard3",  "46.1%", "150/176",26,173, 20, 31),
        ],
        "insight": "1024-token output cap is the root cause of ~90% of failures. Model needs more tokens to complete the decision procedure.",
        "failures": [
            ("Token-limit truncation \u2192 silent TRUE default", "~288 (78%)",
             "Reasoning chain exhausts 1024-token budget mid-step. No VERDICT emitted; harness defaults to TRUE. hard1: 32/35 confab at limit, hard3: 165/173 at limit."),
            ("'Constant operation' rule misfire", "~13 (4%)",
             "Non-truncated false positive. Checks 'disjoint variable sets' but gets disjointness wrong, or Eq1 doesn't actually force a constant operation."),
            ("'Projection' rule misfire", "~6 (2%)",
             "Verifies only variable-non-overlap condition, ignoring structural requirements. Declares TRUE when preconditions aren't met."),
            ("Token-limit truncation \u2192 silent FALSE default", "~40 (11%)",
             "Mirror of #1: all TRUE\u2192FALSE errors also hit 1024 tokens. Model never reaches TRUE conclusion before cutoff."),
            ("Parse errors (empty response)", "40 (11%)",
             "Empty reasoning + empty content, exactly 1024 tokens. Some are API timeouts (latency >1500s). Not malformed output -- just never produced VERDICT."),
        ],
    },
]

APPENDIX_HEADERS = ["Dataset", "Accuracy", "TRUE\u2192TRUE", "TRUE\u2192FALSE",
                    "FALSE\u2192TRUE", "FALSE\u2192FALSE", "Parse Err"]


def build_appendix_slide(prs, model_info):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG)

    # Title
    add_text_box(slide, Inches(0.4), Inches(0.2), Inches(9.2), Inches(0.45),
                 f"Appendix: {model_info['display']}", font_size=22, color=WHITE, bold=True)
    add_text_box(slide, Inches(0.4), Inches(0.6), Inches(9.2), Inches(0.3),
                 f"v4_10KB_cheatsheet.md  |  Backend: OpenRouter  |  Model: {model_info['name']}",
                 font_size=10, color=LIGHT)

    # Table
    ds_rows = model_info["datasets"]
    rows = len(ds_rows) + 1
    cols = 7
    tbl_left = Inches(0.3)
    tbl_top = Inches(1.1)
    tbl_width = Inches(9.4)
    tbl_height = Inches(0.35 * (rows + 0.5))
    shape = slide.shapes.add_table(rows, cols, tbl_left, tbl_top, tbl_width, tbl_height)
    table = shape.table

    col_widths = [Inches(1.0), Inches(1.2), Inches(1.4), Inches(1.4), Inches(1.4), Inches(1.4), Inches(1.0)]
    for i, w in enumerate(col_widths):
        table.columns[i].width = w

    # Header
    for j, hdr in enumerate(APPENDIX_HEADERS):
        set_cell(table.cell(0, j), hdr, font_size=10, color=WHITE, bold=True, bg_color=HEADER_BG)

    # Data
    for i, (ds_name, acc, tt, tf, ft, ff, pe) in enumerate(ds_rows):
        row_bg = ROW_EVEN if i % 2 == 0 else ROW_ODD
        r = i + 1
        set_cell(table.cell(r, 0), ds_name, font_size=11, color=WHITE, bold=True, align=PP_ALIGN.LEFT, bg_color=row_bg)
        set_cell(table.cell(r, 1), acc, font_size=11, color=acc_color(acc), bold=True, bg_color=row_bg)
        set_cell(table.cell(r, 2), str(tt), font_size=11, color=ACCENT, bold=False, bg_color=row_bg)
        set_cell(table.cell(r, 3), str(tf), font_size=11, color=GOLD if tf > 0 else ACCENT, bold=False, bg_color=row_bg)
        # FALSE->TRUE (confabulation) -- highlight red if significant
        ft_color = RED if ft > 10 else (GOLD if ft > 0 else ACCENT)
        set_cell(table.cell(r, 4), str(ft), font_size=11, color=ft_color, bold=(ft > 10), bg_color=row_bg)
        set_cell(table.cell(r, 5), str(ff), font_size=11, color=ACCENT, bold=False, bg_color=row_bg)
        pe_color = RED if pe > 10 else (GOLD if pe > 0 else ACCENT)
        set_cell(table.cell(r, 6), str(pe), font_size=11, color=pe_color, bold=False, bg_color=row_bg)

    # Insight one-liner
    insight_y = tbl_top + tbl_height + Inches(0.15)
    add_text_box(slide, Inches(0.4), insight_y, Inches(9.2), Inches(0.3),
                 model_info["insight"], font_size=9, color=LIGHT)

    # Failure analysis table
    failures = model_info.get("failures", [])
    if failures:
        fa_top = insight_y + Inches(0.35)
        add_text_box(slide, Inches(0.4), fa_top, Inches(3.0), Inches(0.25),
                     "FAILURE ANALYSIS (from JSON reasoning)", font_size=9, color=ACCENT, bold=True)
        fa_top += Inches(0.28)

        fa_rows = len(failures) + 1
        fa_cols = 3
        fa_shape = slide.shapes.add_table(fa_rows, fa_cols, Inches(0.3), fa_top, Inches(9.4), Inches(0.28 * fa_rows))
        fa_table = fa_shape.table
        fa_table.columns[0].width = Inches(3.2)
        fa_table.columns[1].width = Inches(0.9)
        fa_table.columns[2].width = Inches(5.3)

        set_cell(fa_table.cell(0, 0), "Failure Pattern", font_size=8, color=WHITE, bold=True, bg_color=HEADER_BG, align=PP_ALIGN.LEFT)
        set_cell(fa_table.cell(0, 1), "Count", font_size=8, color=WHITE, bold=True, bg_color=HEADER_BG)
        set_cell(fa_table.cell(0, 2), "Detail", font_size=8, color=WHITE, bold=True, bg_color=HEADER_BG, align=PP_ALIGN.LEFT)

        for fi, (pattern, count, detail) in enumerate(failures):
            r = fi + 1
            row_bg = ROW_EVEN if fi % 2 == 0 else ROW_ODD
            set_cell(fa_table.cell(r, 0), pattern, font_size=7, color=WHITE, bold=False, bg_color=row_bg, align=PP_ALIGN.LEFT)
            set_cell(fa_table.cell(r, 1), count, font_size=7, color=GOLD, bold=True, bg_color=row_bg)
            set_cell(fa_table.cell(r, 2), detail, font_size=7, color=LIGHT, bold=False, bg_color=row_bg, align=PP_ALIGN.LEFT)

    # Footer
    add_text_box(slide, Inches(0.4), Inches(5.2), Inches(9.2), Inches(0.25),
                 "SAIR Mathematics Distillation Challenge -- Equational Theories Stage 1  |  Deadline: 2026-04-20",
                 font_size=8, color=LIGHT, alignment=PP_ALIGN.CENTER)


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------

def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)  # 16:9

    build_slide1(prs)
    build_slide2(prs)
    for model_info in MODEL_DETAILS:
        build_appendix_slide(prs, model_info)

    out = Path(__file__).parent / "v4_hard_sweep_summary.pptx"
    prs.save(str(out))
    print(f"Saved: {out}")


if __name__ == "__main__":
    main()
